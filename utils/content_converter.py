import os
from typing import Dict, Any, Optional
import logging
import re
from bs4 import BeautifulSoup
import requests
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from pytube import YouTube
from models.advanced_ai import ai_model
from models.text_processor import TextProcessor

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ContentConverter:
    def __init__(self):
        self.text_processor = TextProcessor()

    def convert_content(self, file_path: str, content_type: str = None) -> Dict[str, Any]:
        """Convert various content types into standardized format."""
        try:
            if not content_type:
                content_type = self._detect_content_type(file_path)

            # Validate URL if it's a web resource
            if content_type in ['article', 'video']:
                if not self._is_valid_url(file_path):
                    raise ValueError("Invalid URL format")

            content = self._extract_content(file_path, content_type)
            if not content:
                raise ValueError(f"No content could be extracted from {content_type}")

            # Process the extracted content
            processed_content = self.text_processor.process_text(content)

            # Store in vector database for similarity search
            content_id = f"{content_type}_{os.path.basename(file_path)}"
            ai_model.store_content_vector(
                content_id=content_id,
                content=content,
                content_type=content_type
            )

            return {
                "content": processed_content.simplified,
                "topics": processed_content.topics,
                "highlighted_terms": processed_content.highlighted_terms,
                "tags": processed_content.tags,
                "content_type": content_type,
                "content_id": content_id
            }

        except requests.exceptions.RequestException as e:
            logger.error(f"Network error accessing URL: {str(e)}")
            raise ValueError(f"Could not access the URL. Please check your internet connection and try again.")
        except YouTube.exceptions.PytubeError as e:
            logger.error(f"YouTube video processing error: {str(e)}")
            raise ValueError(f"Could not process YouTube video. The video might be private or unavailable.")
        except Exception as e:
            logger.error(f"Error converting content: {str(e)}")
            raise

    def _is_valid_url(self, url: str) -> bool:
        """Validate URL format."""
        url_pattern = re.compile(
            r'^https?://'  # http:// or https://
            r'(?:(?:[A-Z0-9](?:[A-Z0-9-]{0,61}[A-Z0-9])?\.)+[A-Z]{2,6}\.?|'  # domain...
            r'localhost|'  # localhost...
            r'\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3})'  # ...or ip
            r'(?::\d+)?'  # optional port
            r'(?:/?|[/?]\S+)$', re.IGNORECASE)
        return url_pattern.match(url) is not None

    def _detect_content_type(self, file_path: str) -> str:
        """Detect content type from file extension or URL."""
        if file_path.startswith(('http://', 'https://')):
            if 'youtube.com' in file_path or 'youtu.be' in file_path:
                return 'video'
            return 'article'

        ext = file_path.lower().split('.')[-1]
        content_types = {
            'pdf': 'pdf',
            'docx': 'document',
            'doc': 'document',
            'pptx': 'presentation',
            'ppt': 'presentation',
            'txt': 'text'
        }
        return content_types.get(ext, 'unknown')

    def _extract_content(self, file_path: str, content_type: str) -> str:
        """Extract text content from various file types."""
        try:
            if content_type == 'pdf':
                return self._extract_from_pdf(file_path)
            elif content_type == 'document':
                return self._extract_from_docx(file_path)
            elif content_type == 'presentation':
                return self._extract_from_pptx(file_path)
            elif content_type == 'article':
                return self._extract_from_url(file_path)
            elif content_type == 'video':
                return self._extract_from_video(file_path)
            elif content_type == 'text':
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            else:
                raise ValueError(f"Unsupported content type: {content_type}")

        except Exception as e:
            logger.error(f"Error extracting content: {str(e)}")
            raise

    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF files."""
        try:
            text_content = []
            with open(file_path, 'rb') as file:
                pdf = PdfReader(file)
                for page in pdf.pages:
                    text_content.append(page.extract_text())
            return '\n'.join(text_content)
        except Exception as e:
            raise ValueError(f"Error reading PDF file: {str(e)}")

    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from Word documents."""
        try:
            doc = Document(file_path)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            raise ValueError(f"Error reading Word document: {str(e)}")

    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations."""
        try:
            pptx = Presentation(file_path)
            text_content = []

            for slide in pptx.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                text_content.append('\n'.join(slide_text))

            return '\n\n'.join(text_content)
        except Exception as e:
            raise ValueError(f"Error reading PowerPoint presentation: {str(e)}")

    def _extract_from_url(self, url: str) -> str:
        """Extract text from web articles."""
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            response = requests.get(url, headers=headers, timeout=10)
            response.raise_for_status()

            soup = BeautifulSoup(response.text, 'html.parser')

            # Remove script and style elements
            for element in soup(['script', 'style', 'nav', 'header', 'footer', 'aside']):
                element.decompose()

            # Get main content
            main_content = soup.find('main') or soup.find('article') or soup.find('body')

            if not main_content:
                raise ValueError("Could not find main content in the article")

            # Get text content
            text = main_content.get_text()

            # Break into lines and remove leading/trailing space
            lines = (line.strip() for line in text.splitlines())

            # Break multi-headlines into a line each
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))

            # Drop blank lines and join
            content = '\n'.join(chunk for chunk in chunks if chunk)

            if not content:
                raise ValueError("No text content found in the article")

            return content

        except requests.exceptions.RequestException as e:
            raise ValueError(f"Error accessing URL: {str(e)}")
        except Exception as e:
            raise ValueError(f"Error processing web article: {str(e)}")

    async def _extract_from_video(self, url: str) -> str:
        """Extract text content from video using AssemblyAI transcription."""
        try:
            # For YouTube videos, get the audio URL first
            if 'youtube.com' in url or 'youtu.be' in url:
                try:
                    yt = YouTube(url)
                    audio_stream = yt.streams.filter(only_audio=True).first()
                    if not audio_stream:
                        raise ValueError("No audio stream found in the video")
                    url = audio_stream.url
                except Exception as e:
                    raise ValueError(f"Error processing YouTube video: {str(e)}")

            # Use advanced AI model's audio processing
            audio_result = await ai_model.process_audio_content(url)
            if not audio_result or 'text' not in audio_result:
                raise ValueError("No transcription generated from the video")

            return audio_result.get('text', '')

        except Exception as e:
            logger.error(f"Error extracting video content: {str(e)}")
            raise

# Create singleton instance
content_converter = ContentConverter()